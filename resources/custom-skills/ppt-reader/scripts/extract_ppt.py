#!/usr/bin/env python3
"""
PPT/PPSX 文件内容提取工具
支持 .pptx 和 .ppsx 格式的演示文稿
"""

import sys
import json
from pathlib import Path

def extract_ppt(ppt_path: str) -> dict:
    """提取 PPT 文件内容"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {
            "success": False,
            "error": "需要安装 python-pptx 库，运行: uv pip install python-pptx"
        }
    
    path = Path(ppt_path)
    if not path.exists():
        return {
            "success": False,
            "error": f"文件不存在: {ppt_path}"
        }
    
    if path.suffix.lower() not in ['.pptx', '.ppsx']:
        return {
            "success": False,
            "error": f"不支持的文件格式: {path.suffix}，仅支持 .pptx 和 .ppsx"
        }
    
    try:
        prs = Presentation(str(path))
        
        result = {
            "success": True,
            "file": str(path),
            "filename": path.name,
            "slide_count": len(prs.slides),
            "slides": [],
            "metadata": {}
        }
        
        # 提取元数据
        core_props = prs.core_properties
        if core_props.title:
            result["metadata"]["title"] = core_props.title
        if core_props.author:
            result["metadata"]["author"] = core_props.author
        if core_props.subject:
            result["metadata"]["subject"] = core_props.subject
        if core_props.created:
            result["metadata"]["created"] = str(core_props.created)
        if core_props.modified:
            result["metadata"]["modified"] = str(core_props.modified)
        
        # 提取每页幻灯片内容
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": i,
                "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "shapes": [],
                "text_content": []
            }
            
            # 提取所有形状
            for shape in slide.shapes:
                shape_info = {
                    "type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                    "name": shape.name
                }
                
                # 提取文本
                if hasattr(shape, "text") and shape.text.strip():
                    shape_info["text"] = shape.text.strip()
                    slide_data["text_content"].append(shape.text.strip())
                
                # 提取图片信息
                if shape.shape_type.name == 'PICTURE':
                    try:
                        image = shape.image
                        shape_info["image"] = {
                            "content_type": image.content_type,
                            "size": f"{image.size[0]}x{image.size[1]}"
                        }
                    except:
                        pass
                
                # 提取表格
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    shape_info["table"] = table_data
                
                slide_data["shapes"].append(shape_info)
            
            # 合并该页所有文本
            slide_data["full_text"] = "\n".join(slide_data["text_content"])
            
            result["slides"].append(slide_data)
        
        # 生成摘要
        all_text = []
        for slide in result["slides"]:
            if slide["full_text"]:
                all_text.append(f"=== 第 {slide['slide_number']} 页 ===\n{slide['full_text']}")
        
        result["full_content"] = "\n\n".join(all_text)
        
        return result
        
    except Exception as e:
        return {
            "success": False,
            "error": f"解析 PPT 文件失败: {str(e)}"
        }


def main():
    # 设置标准输出编码为 UTF-8
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    
    if len(sys.argv) < 2:
        print("用法: python extract_ppt.py <ppt文件路径>")
        print("支持格式: .pptx, .ppsx")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    result = extract_ppt(ppt_path)
    
    # 输出 JSON 结果
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
